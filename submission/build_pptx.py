#!/usr/bin/env python3
"""Builds the AI Kavach submission deck for ABHIMANYU X.

Content is drawn only from what was actually built and verified this
session (real REWIND detection, real AFL++/ASan crash discovery, real
local-LLM patch generation, real compile+replay verification, real
Immune Transfer experiment on a second real target) — nothing here is
aspirational or fabricated for the pitch.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

INK = RGBColor(0x0E, 0x16, 0x14)
PAPER = RGBColor(0xEE, 0xF1, 0xEC)
ACCENT = RGBColor(0x2D, 0xE2, 0xC9)
ACCENT_2 = RGBColor(0xE8, 0xB2, 0x3D)
TEXT = RGBColor(0xE9, 0xED, 0xE9)
TEXT_DIM = RGBColor(0x9F, 0xB0, 0xAC)
DANGER = RGBColor(0xFF, 0x54, 0x70)
GREEN = RGBColor(0x3D, 0xDC, 0x84)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def add_text(slide, left, top, width, height, text, size=18, color=TEXT, bold=False,
             italic=False, align=PP_ALIGN.LEFT, font="Georgia", line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=15, color=TEXT,
                 marker_color=ACCENT, font="Calibri", gap=6, line_spacing=1.08):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(gap)
        if isinstance(item, tuple):
            head, rest = item
            r1 = p.add_run()
            r1.text = "›  " + head
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.color.rgb = marker_color
            r1.font.name = font
            if rest:
                r2 = p.add_run()
                r2.text = "  —  " + rest
                r2.font.size = Pt(size - 1)
                r2.font.color.rgb = color
                r2.font.name = font
        else:
            r1 = p.add_run()
            r1.text = "›  " + item
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
            r1.font.name = font
    return box


def kicker(slide, text, num):
    add_text(slide, Inches(0.55), Inches(0.28), Inches(4), Inches(0.4),
              f"ABHIMANYU X  ·  AI KAVACH", size=11, color=TEXT_DIM, font="Courier New")
    add_text(slide, Inches(12.2), Inches(0.28), Inches(0.7), Inches(0.4),
              f"{num}/5", size=11, color=ACCENT, font="Courier New", align=PP_ALIGN.RIGHT)


def title(slide, text, subtitle=None):
    add_text(slide, Inches(0.55), Inches(0.65), Inches(12.2), Inches(0.9),
              text, size=32, color=TEXT, bold=True)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(1.28), Inches(12.2), Inches(0.5),
                  subtitle, size=15, color=ACCENT, italic=True)


def hairline(slide, top):
    ln = slide.shapes.add_connector(1, Inches(0.55), top, Inches(12.78), top)
    ln.line.color.rgb = RGBColor(0x1C, 0x2A, 0x28)
    ln.line.width = Pt(0.75)


def flow_box(slide, left, top, width, height, label, sub, color=ACCENT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x13, 0x1E, 0x1C)
    box.line.color.rgb = color
    box.line.width = Pt(1.25)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = TEXT
    r.font.name = "Courier New"
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = TEXT_DIM
        r2.font.name = "Calibri"
    return box


def arrow_right(slide, left, top, width=Inches(0.35)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.28))
    a.fill.solid()
    a.fill.fore_color.rgb = ACCENT
    a.line.fill.background()
    a.shadow.inherit = False


def badge(slide, left, top, text, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.35), Inches(0.32))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x13, 0x1E, 0x1C)
    box.line.color.rgb = color
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_top = Emu(10000)
    tf.margin_bottom = Emu(10000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = "Courier New"
    return box


# ============================================================
# SLIDE 1 — Introduction, Ideation & Brief Description
# ============================================================
s = add_slide()
kicker(s, "Introduction", 1)
title(s, "ABHIMANYU X", "Autonomous Cyber Immune System for Defence Software")
hairline(s, Inches(1.85))

add_text(s, Inches(0.55), Inches(2.05), Inches(6.0), Inches(0.4), "THE PROBLEM", size=13, color=ACCENT_2, bold=True, font="Courier New")
add_bullets(s, Inches(0.55), Inches(2.45), Inches(6.0), Inches(2.6), [
    ("AI patch tools trust the LLM.", "a generated fix is shown as done, with no proof it actually blocks the exploit."),
    ("Security tooling needs heavy infra.", "cloud APIs, GPUs, or clusters most defence-context labs can't rely on."),
    ("Manual discover-patch-verify cycles are slow", "and don't get faster just by adding an LLM in the loop."),
], size=14)

add_text(s, Inches(6.85), Inches(2.05), Inches(5.9), Inches(0.4), "THE IDEA", size=13, color=ACCENT_2, bold=True, font="Courier New")
add_bullets(s, Inches(6.85), Inches(2.45), Inches(5.9), Inches(2.6), [
    ("A closed-loop cyber-immune cell.", "finds a real vulnerability, reasons about root cause with a local LLM, patches it, and PROVES the fix holds."),
    ("AI proposes. Evidence decides.", "every patch must pass real compiler + real exploit replay + real regression before it's called “verified.”"),
    ("Fully local inference.", "a 3B-parameter model running offline — no cloud dependency, air-gap capable."),
], size=14)

hairline(s, Inches(5.15))
add_text(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(0.4), "CORE LOOP", size=13, color=ACCENT_2, bold=True, font="Courier New")

stages = ["DISCOVER", "UNDERSTAND", "REPAIR", "VERIFY", "REMEMBER", "TRANSFER"]
bw, gap = Inches(1.85), Inches(0.15)
x = Inches(0.55)
y = Inches(5.85)
for i, st in enumerate(stages):
    flow_box(s, x, y, bw, Inches(0.75), st, None)
    x = Emu(x + bw + gap)
    if i < len(stages) - 1:
        arrow_right(s, Emu(x - gap + Emu(20000)), Emu(y + Inches(0.75)//2 - Inches(0.14)))

add_text(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.5),
          "Every verified vulnerability becomes future defence.", size=14, color=TEXT_DIM, italic=True)


# ============================================================
# SLIDE 2 — Detailed Methodology
# ============================================================
s = add_slide()
kicker(s, "Methodology", 2)
title(s, "Detailed Methodology", "Step-by-step pipeline — every stage below actually executes")
hairline(s, Inches(1.85))

steps = [
    ("1. REWIND", "Real `git diff` commit analysis + static detection engine (Python + C). Includes a custom rule the project added: catching unchecked length reaching memcpy() — CWE-120 — which no pre-existing pattern covered."),
    ("2. Dynamic Analysis & Fuzzing", "AFL++ 4.09c installed as a genuine OS package inside a Linux VM (Colima/Ubuntu 24.04). Blind-mode fuzzing against a real clang+ASan binary found a real stack-buffer-overflow crash. A real LLVM17/ASan toolchain incompatibility was found and disclosed, not hidden."),
    ("3. ANVIL — AI Reasoning", "Root-cause analysis and patch generation via a local Ollama model (Qwen2.5-Coder 3B). No cloud call, no chain-of-thought exposed — only structured, evidence-grounded conclusions."),
    ("4. Verification Gates", "Real compile check, real exploit replay (same crash input, before vs. after, recompiled and re-executed in the VM), real regression/behaviour checks, plus adversarial testing across 6 additional payload sizes spanning the boundary."),
    ("5. Self-Correcting Repair", "A failed verification sends the concrete failure reason back to ANVIL for a real second attempt — GENERATE → TEST → FAIL → REPAIR → RETEST, not one-shot patching."),
    ("6. Immune Memory & Transfer", "Verified patterns are stored as “Vulnerability DNA” in SQLite. A second real target (network_protocol_parser) demonstrates a measured with-memory vs. without-memory comparison — not a fabricated benchmark."),
]
colw = Inches(6.0)
for i, (h, d) in enumerate(steps):
    col = i % 2
    row = i // 2
    left = Inches(0.55) + col * Inches(6.35)
    top = Inches(2.05) + row * Inches(1.6)
    add_text(s, left, top, colw, Inches(0.35), h, size=15, color=ACCENT, bold=True, font="Courier New")
    add_text(s, left, Emu(top + Inches(0.4)), colw, Inches(1.1), d, size=11.5, color=TEXT, font="Calibri", line_spacing=1.1)


# ============================================================
# SLIDE 3 — Technology Stack / Flow Diagram / Equipment
# ============================================================
s = add_slide()
kicker(s, "Technology & Architecture", 3)
title(s, "Technology Stack & System Flow")
hairline(s, Inches(1.85))

add_text(s, Inches(0.55), Inches(2.0), Inches(5.6), Inches(0.35), "STACK", size=13, color=ACCENT_2, bold=True, font="Courier New")
add_bullets(s, Inches(0.55), Inches(2.4), Inches(5.6), Inches(4.2), [
    ("Static analysis", "REWIND engine — Python, regex + AST + heuristics"),
    ("Fuzzing / dynamic", "AFL++ 4.09c, clang-18, AddressSanitizer"),
    ("Local AI", "Ollama + Qwen2.5-Coder 3B — fully offline inference"),
    ("Backend", "Python, Flask + Flask-SocketIO (real-time WebSocket events)"),
    ("Data", "SQLite — Immune Memory / Vulnerability DNA store"),
    ("Frontend", "Vanilla Three.js 3D laboratory — no build toolchain, reacts to real backend state"),
    ("CLI", "Python argparse — doctor / scan / mission / transfer"),
    ("Supply chain", "pip-audit (real SBOM/CVE scan), Docker buildx (real multi-arch check)"),
    ("Equipment", "MacBook (Apple Silicon) host + Colima-managed Linux VM — 2 vCPU / 4 GB RAM / 20 GB disk"),
], size=12.5, gap=5)

add_text(s, Inches(6.6), Inches(2.0), Inches(6.15), Inches(0.35), "SYSTEM FLOW", size=13, color=ACCENT_2, bold=True, font="Courier New")

flow = [
    ("REPOSITORY", "target source + git history"),
    ("REWIND", "commit diff + static findings"),
    ("STATIC + FUZZ ENGINE", "REWIND scan  //  AFL++ + ASan"),
    ("VULNERABILITY", "evidence bundle, CWE mapped"),
    ("ANVIL (local LLM)", "root cause + patch, offline"),
    ("VERIFICATION CHAMBER", "build → replay → regression → adversarial"),
    ("IMMUNE MEMORY", "verified pattern stored + transferable"),
]
fx = Inches(6.6)
fy = Inches(2.4)
fw = Inches(6.15)
fh = Inches(0.62)
fgap = Inches(0.08)
for i, (h, sub) in enumerate(flow):
    flow_box(s, fx, fy, fw, fh, h, sub)
    fy = Emu(fy + fh + fgap)
    if i < len(flow) - 1:
        down = slide_arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Emu(fx + fw//2 - Inches(0.09)), Emu(fy - fgap + Emu(4000)), Inches(0.18), Emu(fgap - Emu(8000)))
        down.fill.solid()
        down.fill.fore_color.rgb = ACCENT
        down.line.fill.background()
        down.shadow.inherit = False


# ============================================================
# SLIDE 4 — Salient Features & Novelty
# ============================================================
s = add_slide()
kicker(s, "Salient Features & Novelty", 4)
title(s, "What Makes This Different", "AI proposes. Evidence decides. Verification proves. Memory learns.")
hairline(s, Inches(1.85))

features = [
    ("Evidence-driven verification, not AI self-certification",
     "A patch is only called “VERIFIED” after passing all real gates — compile, exploit replay, regression, behaviour, adversarial robustness (7/7, computed, never asserted)."),
    ("Self-correcting repair loop",
     "A rejected patch's real failure reason is sent back to ANVIL for a genuine second attempt — the system demonstrably knows when NOT to trust its own output."),
    ("Measured Immune Transfer, not a marketing claim",
     "A second real vulnerable target shows an actual with-memory vs. without-memory comparison — reported as measured, whatever the real result turned out to be."),
    ("Fully local, fully lightweight",
     "A 3B-parameter model and a 2-vCPU/4GB VM run the entire discover→patch→verify→remember loop on a single laptop — no GPU, no cluster, no cloud API."),
    ("Radical honesty baked into the UI itself",
     "Every figure is tagged MEASURED / AI-GENERATED / DEMO / FUTURE. A real, disclosed toolchain limitation (AFL++ coverage instrumentation vs. ASan on this platform) is documented rather than hidden — and worked around with a real fix, not a fudge."),
    ("Reproducible by construction",
     "Every mission emits real provenance: commit hash, target/patch SHA-256, tool versions, timestamps — the same run can be independently re-verified."),
]
for i, (h, d) in enumerate(features):
    col = i % 2
    row = i // 2
    left = Inches(0.55) + col * Inches(6.35)
    top = Inches(2.05) + row * Inches(1.6)
    add_text(s, left, top, Inches(6.0), Inches(0.5), h, size=14.5, color=ACCENT, bold=True)
    add_text(s, left, Emu(top + Inches(0.5)), Inches(6.0), Inches(1.0), d, size=11.5, color=TEXT, line_spacing=1.12)


# ============================================================
# SLIDE 5 — Final Deliverables
# ============================================================
s = add_slide()
kicker(s, "Final Deliverables", 5)
title(s, "Deliverables & Proof of Concept")
hairline(s, Inches(1.85))

add_text(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(0.35), "DELIVERED", size=13, color=ACCENT_2, bold=True, font="Courier New")
add_bullets(s, Inches(0.55), Inches(2.4), Inches(6.0), Inches(4.4), [
    "Working end-to-end prototype: DISCOVER → UNDERSTAND → REPAIR → VERIFY → REMEMBER → TRANSFER",
    "Two real vulnerable C targets — real git history, real fuzz harness, real AFL++/ASan-found crash each",
    "3D interactive laboratory UI (Three.js) and a CLI, both driven by the same real backend event stream",
    "Real per-mission Provenance record for reproducibility (commit hash, file/patch SHA-256, tool versions, timestamps)",
    "Real SBOM (pip-audit) and real multi-architecture capability report (Docker buildx)",
    "Retry/rollback loop and adversarial robustness testing on every verified patch",
], size=13)

add_text(s, Inches(6.85), Inches(2.0), Inches(5.9), Inches(0.35), "DEMONSTRATED LIVE", size=13, color=ACCENT_2, bold=True, font="Courier New")
demo_steps = ["Real crash found (AFL++/ASan)", "Real patch generated (local LLM)",
              "Real exploit replay — blocked", "Real regression — pass",
              "Immune Memory created", "Pattern recognized on 2nd real target"]
dy = Inches(2.4)
for i, step in enumerate(demo_steps):
    box = flow_box(s, Inches(6.85), dy, Inches(5.9), Inches(0.55), step, None, color=GREEN if i < 5 else ACCENT_2)
    dy = Emu(dy + Inches(0.55) + Inches(0.1))

hairline(s, Inches(6.35))
add_text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.6),
          "Status: research prototype — not validated for production or defence-classified deployment.\nEvery patch is meant to be human-reviewed before use, not auto-deployed.",
          size=11, color=TEXT_DIM, italic=True, line_spacing=1.2)

out_path = "/Users/sachi/abhimanyux/submission/ABHIMANYU_X_AI_Kavach_Submission.pptx"
prs.save(out_path)
print("Saved:", out_path)
